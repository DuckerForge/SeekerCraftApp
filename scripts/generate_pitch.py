"""
SeekerCraft Hackathon Pitch Deck Generator
Run: python3 scripts/generate_pitch.py
Output: scripts/seekercraft-pitch.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE = "/home/oliver/SeekerCraft"
ASSETS = f"{BASE}/assets/images"
OUT = f"{BASE}/scripts/seekercraft-pitch.pptx"

# ── Palette ───────────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x00, 0x20)
C_CARD      = RGBColor(0x16, 0x00, 0x35)
C_CARD2     = RGBColor(0x1E, 0x00, 0x42)
C_PURPLE    = RGBColor(0x7B, 0x2F, 0xFF)
C_BLUE      = RGBColor(0x1E, 0x90, 0xFF)
C_GOLD      = RGBColor(0xFF, 0xD7, 0x00)
C_GREEN     = RGBColor(0x00, 0xFF, 0x88)
C_ORANGE    = RGBColor(0xFF, 0x66, 0x00)
C_PINK      = RGBColor(0xFF, 0x45, 0xAA)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xCC, 0xCC, 0xFF)
C_HEADER_BG = RGBColor(0x1A, 0x00, 0x3A)

# ── Helpers ───────────────────────────────────────────────────────────────────

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, color, line=False):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line:
        s.line.color.rgb = color
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=16, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = "Segoe UI" if os.name == "nt" else "DejaVu Sans"
    return tb

def img(slide, path, x, y, w, h=None):
    if not path or not os.path.exists(path):
        return None
    try:
        if h:
            return slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
        return slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w))
    except Exception:
        return None

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.05, C_HEADER_BG)
    rect(slide, 0, 1.05, 13.33, 0.05, C_GOLD)
    txt(slide, title, 0.4, 0.08, 11, 0.65, size=34, bold=True, color=C_WHITE)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.7, 11.5, 0.38, size=14, color=C_LIGHT)

def tag_bar(slide, text, y=7.1):
    rect(slide, 0, y, 13.33, 0.4, C_HEADER_BG)
    txt(slide, text, 0, y + 0.05, 13.33, 0.3, size=13, color=C_LIGHT,
        align=PP_ALIGN.CENTER)

# ── Presentation setup ────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)

img(s, f"{ASSETS}/banner_1200x600.png", 0, 0, 13.33, 4.3)
rect(s, 0, 4.3, 13.33, 3.2, C_BG)
rect(s, 0, 4.3, 13.33, 0.05, C_GOLD)

img(s, f"{ASSETS}/NewLogo.png", 0.35, 4.55, 1.85)
txt(s, "SEEKERCRAFT",         2.55, 4.38, 10.5, 1.05, size=52, bold=True,  color=C_WHITE)
txt(s, "Shoot. Build. Conquer.",          2.55, 5.32, 10.5, 0.65, size=26,              color=C_GOLD)
txt(s, "A Peggle-style physics game built natively for Solana Mobile\n"
       "with full SKR token integration & on-chain payments.",
       2.55, 5.9, 10.2, 0.85, size=16, color=C_LIGHT)

tag_bar(s, "Solana Mobile  |  Expo / React Native  |  Firebase Realtime DB  |  SKR Token  |  60fps Worklet Physics")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — THE PROBLEM
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
header(s, "THE PROBLEM", "Mobile gaming on Solana is missing something real")

problems = [
    ("Games are built for devs, not players",    "Most crypto games funnel rewards back to the dev.\nThe community that plays and creates sees little benefit.",  C_ORANGE),
    ("No tools for community content",           "Mobile crypto gaming has no real way for players\nto create, share, and monetize their own levels.",            C_PURPLE),
    ("No real player-to-player competition",     "No ranked duels, no shared challenges, no social layer.\nPlayers exist in isolation with nothing to compete for.", C_BLUE),
    ("SKR token has no gaming home",             "No fun, casual way to use SKR on Seeker.\nToken utility stays theoretical rather than playful.",               C_GOLD),
]

for i, (title, desc, color) in enumerate(problems):
    col, row = i % 2, i // 2
    x, y, w, h = 0.25 + col * 6.55, 1.25 + row * 2.85, 6.25, 2.6
    rect(s, x, y, w, h, C_CARD)
    rect(s, x, y, 0.07, h, color)
    txt(s, title, x + 0.22, y + 0.18, w - 0.35, 0.55, size=19, bold=True, color=color)
    txt(s, desc,  x + 0.22, y + 0.85, w - 0.35, 1.55, size=14, color=C_LIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — THE SOLUTION
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
header(s, "THE SOLUTION", "SeekerCraft: every interaction is Solana-native")

pillars = [
    ("PLAY",    "Addictive Peggle-style\nball physics",        C_BLUE,   f"{ASSETS}/Peg/peg_gold.png"),
    ("BUILD",   "Create & publish levels\nfor the community",  C_PURPLE, f"{ASSETS}/Peg/bump.png"),
    ("COMPETE", "PvP duels with\non-chain score submission",   C_ORANGE, f"{ASSETS}/Icons/duel.png"),
    ("EARN",    "160+ achievements\nwith on-chain unlocks",    C_GOLD,   f"{ASSETS}/Icons/coppa.png"),
    ("TIP",     "Reward creators\nwith SKR tokens (90%)",      C_GREEN,  f"{ASSETS}/Peg/skr.png"),
]

cw = 13.33 / 5
for i, (label, desc, color, icon_path) in enumerate(pillars):
    x, y, w, h = i * cw + 0.1, 1.2, cw - 0.2, 6.05
    rect(s, x, y, w, h, C_CARD)
    rect(s, x, y, w, 0.07, color)
    img(s, icon_path, x + w / 2 - 0.35, y + 0.3, 0.7)
    txt(s, label, x, y + 1.95, w, 0.6, size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, desc,  x + 0.1, y + 2.65, w - 0.2, 1.6, size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — HOW IT WORKS
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
header(s, "HOW IT WORKS", "Simple to pick up, deep to master")

steps = [
    ("1", "Connect Wallet",   "Tap to connect your Solana wallet\nvia Mobile Wallet Adapter"),
    ("2", "Aim & Shoot",      "Drag to aim, release to fire\nyour ball into the peg grid"),
    ("3", "Hit Gold Pegs",    "Clear all golden SKR pegs\nto win the level"),
    ("4", "Use Power-Ups",    "Multiball, Fireball, Slow Bucket,\nLong Shot, Vortex"),
    ("5", "Epic Slow-Mo",     "Cinematic zoom on the last\nSKR peg, Peggle 2 style"),
    ("6", "Score & Rank",     "Build combos up to 10x,\nclimb the global leaderboard"),
]

for i, (num, title, desc) in enumerate(steps):
    col, row = i % 3, i // 3
    x, y, w, h = 0.25 + col * 4.35, 1.25 + row * 2.95, 4.08, 2.65
    rect(s, x, y, w, h, C_CARD)
    rect(s, x, y, 0.6, 0.6, C_PURPLE)
    txt(s, num,   x, y, 0.6, 0.6, size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, x + 0.72, y + 0.08, w - 0.85, 0.5, size=18, bold=True, color=C_WHITE)
    txt(s, desc,  x + 0.18, y + 0.82, w - 0.3, 1.65, size=13, color=C_LIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — KEY FEATURES
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
header(s, "KEY FEATURES")

features = [
    ("LEVEL EDITOR", C_PURPLE, f"{ASSETS}/Peg/bump.png",
     ["12x16 grid with 7 unique peg types",
      "Curves, vortexes, teleporters, moving walls",
      "10 SKR upload fee to reduce spam",
      "Cloud sync: edit offline, publish anytime"]),
    ("COMMUNITY", C_BLUE, f"{ASSETS}/Icons/players.png",
     ["Browse & play other players' levels",
      "Rate levels 1-5 stars with live feed",
      "Global leaderboard (all-time + weekly)",
      "Real-time activity & donation feed"]),
    ("PVP DUELS", C_ORANGE, f"{ASSETS}/Icons/duel.png",
     ["Challenge any player on the leaderboard",
      "Both play same level, highest score wins",
      "On-chain score submission via Firebase",
      "Win/Draw/Loss record tracked per player"]),
    ("160+ ACHIEVEMENTS", C_GOLD, f"{ASSETS}/Icons/coppa.png",
     ["7 categories: Player, Score, Creator, Duel...",
      "On-chain unlock for 0.50 SOL (permanent)",
      "Rarity tiers: Common, Rare, Epic, Legendary",
      "Donation & Social milestones included"]),
]

for i, (title, color, icon_path, points) in enumerate(features):
    col, row = i % 2, i // 2
    x, y, w, h = 0.2 + col * 6.65, 1.2 + row * 2.98, 6.35, 2.72
    rect(s, x, y, w, h, C_CARD)
    rect(s, x, y, 0.07, h, color)
    img(s, icon_path, x + 0.22, y + 0.22, 0.5)
    txt(s, title, x + 1.08, y + 0.2, w - 1.22, 0.52, size=18, bold=True, color=color)
    for j, pt in enumerate(points):
        txt(s, f"• {pt}", x + 1.08, y + 0.82 + j * 0.46, w - 1.22, 0.42, size=12.5, color=C_LIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — SKR TOKEN INTEGRATION
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
header(s, "SKR TOKEN INTEGRATION", "SeekerCraft deeply integrates the SKR (Seeker) token")

img(s, f"{ASSETS}/Peg/skr.png",  5.9, 1.12, 1.55)
img(s, f"{ASSETS}/Peg/sol.png",  7.7, 1.12, 1.55)

cols_data = [
    ("DONATIONS", C_GREEN, [
        "Tip level creators with SKR tokens",
        "90% goes to creator / 10% dev fund",
        "Preset amounts: $0.50 / $5 / $15",
        "Custom SKR amount supported",
    ]),
    ("ON-CHAIN TX", C_BLUE, [
        "SPL token transfers via MWA",
        "Auto-creates Associated Token Accounts",
        "Real-time SKR balance checks",
        "Full TX preview before signing",
    ]),
    ("GAME REWARDS", C_GOLD, [
        "Golden SKR pegs = win condition",
        "SKR balance shown in rankings",
        "Donation milestones unlock achievements",
        "10 SKR upload fee to reduce spam",
    ]),
]

cw = 13.33 / 3
for i, (title, color, items) in enumerate(cols_data):
    x, y, w, h = i * cw + 0.15, 1.35, cw - 0.3, 5.9
    rect(s, x, y, w, h, C_CARD)
    rect(s, x, y, w, 0.07, color)
    txt(s, title, x, y + 0.15, w, 0.6, size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
    rect(s, x + 0.35, y + 0.88, w - 0.7, 0.03, color)
    for j, item in enumerate(items):
        txt(s, f"• {item}", x + 0.2, y + 1.05 + j * 1.12, w - 0.4, 1.0, size=14, color=C_LIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — TECH STACK
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
header(s, "TECH STACK", "Mobile-native, production-ready architecture")

tech = [
    ("Frontend",    "Expo / React Native (expo-router)",                              C_BLUE),
    ("Rendering",   "@shopify/react-native-skia, Canvas GPU rendering",               C_PURPLE),
    ("Physics",     "react-native-reanimated worklets, 60fps, no JS bridge",          C_PURPLE),
    ("Wallet",      "@solana-mobile/mobile-wallet-adapter-protocol-web3js",           RGBColor(0x9C, 0x45, 0xFF)),
    ("Blockchain",  "@solana/web3.js + @solana/spl-token, SOL and SPL transfers",     C_GREEN),
    ("Backend",     "Firebase Realtime Database: levels, scores, activity, duels",    RGBColor(0xFF, 0xA5, 0x00)),
    ("Audio",       "expo-audio: 5 music tracks + 8 SFX, fever music system",         C_BLUE),
    ("Security",    "dotenv + expo-constants, zero hardcoded keys in source",         C_GREEN),
]

for i, (layer, detail, color) in enumerate(tech):
    col, row = i % 2, i // 2
    x, y, w, h = 0.25 + col * 6.6, 1.25 + row * 1.42, 6.3, 1.3
    rect(s, x, y, w, h, C_CARD)
    rect(s, x, y, 0.07, h, color)
    txt(s, layer.upper(), x + 0.2, y + 0.1,  1.4, 0.45, size=12, bold=True, color=color)
    txt(s, detail,         x + 0.2, y + 0.58, w - 0.35, 0.58, size=15, color=C_WHITE)

rect(s, 0.25, 7.0, 12.83, 0.35, C_HEADER_BG)
txt(s, "All on-chain: SOL transfers  |  SKR SPL transfers  |  ComputeBudgetProgram priority fees  |  Full MWA authorize / sign / confirm",
    0.25, 7.04, 12.83, 0.28, size=12, color=C_GOLD, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — MOBILE-FIRST DESIGN
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
header(s, "MOBILE-FIRST DESIGN", "Built exclusively for Seeker / Solana Mobile")

left_feats = [
    ("Touch Controls",   "Drag to aim, release to fire. Zero buttons needed.",          C_BLUE),
    ("Haptic Feedback",  "Every peg hit, combo, achievement triggers vibration",        C_PURPLE),
    ("Portrait UI",      "Candy Crush style, fully optimized for one-hand play",        C_BLUE),
    ("Cinematic Slow-Mo","Peggle 2-style zoom on the last golden SKR peg",              C_PURPLE),
]
right_feats = [
    ("7 Languages",     "EN / FR / DE / ES / ZH / JA / KO, fully localized",          C_GOLD),
    ("Neon Visual FX",  "Sparkle particles, screen shake, combo overlays, confetti",   C_GREEN),
    ("Offline + Cloud", "Edit levels offline, auto-sync to Firebase on publish",        C_GOLD),
    ("4 Music Tracks",  "Alternating BGM with fever mode & ding combo system",          C_GREEN),
]

for i, (title, desc, color) in enumerate(left_feats):
    y = 1.25 + i * 1.48
    rect(s, 0.25, y, 5.95, 1.33, C_CARD)
    rect(s, 0.25, y, 0.07, 1.33, color)
    txt(s, title, 0.45, y + 0.1,  5.55, 0.48, size=17, bold=True, color=color)
    txt(s, desc,  0.45, y + 0.62, 5.55, 0.6,  size=13, color=C_LIGHT)

for i, (title, desc, color) in enumerate(right_feats):
    y = 1.25 + i * 1.48
    rect(s, 7.13, y, 5.95, 1.33, C_CARD)
    rect(s, 7.13, y, 0.07, 1.33, color)
    txt(s, title, 7.33, y + 0.1,  5.55, 0.48, size=17, bold=True, color=color)
    txt(s, desc,  7.33, y + 0.62, 5.55, 0.6,  size=13, color=C_LIGHT)

rect(s, 6.55, 1.25, 0.05, 5.95, RGBColor(0x30, 0x00, 0x55))

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — TRACTION & NUMBERS
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
header(s, "TRACTION & NUMBERS", "Built by a solo developer in ~3 months")

stats = [
    ("160+", "Achievements\nacross 7 categories",         C_GOLD),
    ("7",    "Peg types + curves\n+ teleporters + walls",  C_PURPLE),
    ("5",    "Power-up types:\nMultiball, Fireball...",    C_BLUE),
    ("8",    "Sound effects +\n5 music tracks",            C_GREEN),
    ("7",    "Languages\nfully supported",                 C_PINK),
    ("1",    "Solo developer\n~3 months to ship",          C_ORANGE),
]

sw = 13.33 / 3
for i, (num, label, color) in enumerate(stats):
    col, row = i % 3, i // 2
    x, y, w, h = 0.2 + col * sw, 1.25 + row * 2.98, sw - 0.3, 2.72
    rect(s, x, y, w, h, C_CARD)
    rect(s, x, y, w, 0.08, color)
    txt(s, num,   x, y + 0.1,  w, 1.4, size=68, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, label, x, y + 1.55, w, 1.0, size=15,             color=C_LIGHT, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — ROADMAP
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
header(s, "ROADMAP")

quarters = [
    ("Q1 2026", "NOW", C_GREEN, [
        "Hackathon submission",
        "Solana dApp Store listing",
        "Community beta launch",
    ]),
    ("Q2 2026", "", C_BLUE, [
        "Tournament mode",
        "Level packs (curated collections)",
        "Enhanced SKR rewards system",
    ]),
    ("Q3 2026", "", C_PURPLE, [
        "Multiplayer real-time mode",
        "NFT level ownership",
        "Cross-platform web version",
    ]),
    ("Q4 2026", "", C_GOLD, [
        "DAO governance",
        "Creator monetization dashboard",
        "Seasonal events & limited achievements",
    ]),
]

# Timeline line
rect(s, 0.5, 2.72, 12.33, 0.07, RGBColor(0x3A, 0x00, 0x6A))

qw = 13.33 / 4
for i, (q, badge, color, items) in enumerate(quarters):
    cx = i * qw + qw / 2 - 0.22
    rect(s, cx, 2.46, 0.44, 0.44, color)  # timeline node

    x, y, w, h = i * qw + 0.15, 1.25, qw - 0.3, 1.1
    rect(s, x, y, w, h, C_HEADER_BG)
    label = f"{q}  [ {badge} ]" if badge else q
    txt(s, label, x, y + 0.2, w, 0.65, size=17, bold=True, color=color, align=PP_ALIGN.CENTER)

    x2, y2, w2, h2 = i * qw + 0.15, 3.05, qw - 0.3, 4.15
    rect(s, x2, y2, w2, h2, C_CARD)
    rect(s, x2, y2, w2, 0.07, color)
    for j, item in enumerate(items):
        txt(s, f"• {item}", x2 + 0.2, y2 + 0.2 + j * 1.22, w2 - 0.35, 1.1, size=14, color=C_LIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — TEAM
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
header(s, "TEAM")

rect(s, 3.3, 1.3, 6.73, 5.75, C_CARD)
rect(s, 3.3, 1.3, 6.73, 0.07, C_PURPLE)

img(s, f"{ASSETS}/NewLogo.png", 5.55, 1.5, 2.23)

txt(s, "DuckerForge", 3.3, 3.95, 6.73, 0.85, size=34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
txt(s, "Solo Developer & Game Designer",
    3.3, 4.75, 6.73, 0.55, size=18, color=C_GOLD, align=PP_ALIGN.CENTER)

rect(s, 4.3, 5.4, 4.73, 0.04, C_PURPLE)

txt(s, "Full-stack mobile developer specializing in\nReact Native, Solana, and game physics.",
    3.3, 5.55, 6.73, 0.85, size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)

txt(s, "github.com/DuckerForge  |  X: @SeekerCraftApp",
    3.3, 6.5, 6.73, 0.45, size=13, color=C_BLUE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — CALL TO ACTION
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)

img(s, f"{ASSETS}/banner_1200x600.png", 0, 0, 13.33, 4.1)
rect(s, 0, 4.1, 13.33, 3.4, C_BG)
rect(s, 0, 4.1, 13.33, 0.06, C_GOLD)

txt(s, "SeekerCraft brings real gaming to Solana Mobile.",
    0, 4.3, 13.33, 0.9, size=30, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
txt(s, "Play. Build. Compete. Earn. Tip.",
    0, 5.1, 13.33, 0.75, size=26, color=C_GOLD, align=PP_ALIGN.CENTER)
txt(s, "All on-chain. All on Seeker.",
    0, 5.78, 13.33, 0.6, size=20, color=C_LIGHT, align=PP_ALIGN.CENTER)

rect(s, 3.5, 6.48, 6.33, 0.04, C_PURPLE)
txt(s, "github.com/DuckerForge/SeekerCraftApp",
    0, 6.62, 13.33, 0.5, size=15, color=C_BLUE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
